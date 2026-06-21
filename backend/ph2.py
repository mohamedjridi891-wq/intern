from pathlib import Path
import pandas as pd
import logging
import zipfile
try:
    import rarfile
except ImportError:
    rarfile = None
import tempfile
import os
import io
import sys
import re
import threading
import warnings
import shutil
from concurrent.futures import ThreadPoolExecutor, as_completed
import psycopg2
from psycopg2.extras import execute_values, RealDictCursor
from dotenv import load_dotenv

load_dotenv()

DB_URL = os.getenv("DB_URL")
logging.basicConfig(
    filename=os.getenv("EXTRACTION_ERRORS", "extraction_errors.log"),
    level=logging.WARNING,
    format="%(asctime)s — %(levelname)s — %(message)s",
)

NUM_WORKERS = int(os.getenv("NUM_WORKERS", 6))
PDF_TIMEOUT_SEC = int(os.getenv("PDF_TIMEOUT_SEC", 15))
OCR_DPI = int(os.getenv("OCR_DPI", 150))
MAX_TEXT_CHARS = int(os.getenv("MAX_TEXT_CHARS", 100_000))
MIN_LANG_CHARS = int(os.getenv("MIN_LANG_CHARS", 50))
OCR_LANG = os.getenv("OCR_LANG", "ara+fra+eng")
REVIEW_FOLDER = os.getenv(
    "REVIEW_FOLDER",
    os.path.join(os.path.dirname(os.path.abspath(__file__)), "review"),
)
ALLOWED_LANGUAGES = {"fr", "en", "ar"}


def _move_to_review(path: str) -> bool:
    if not path:
        return False

    if "::" in path:
        logging.warning(f"Skipping virtual archive path move to review: {path}")
        return False

    path_obj = Path(path)
    if not path_obj.exists():
        logging.warning(f"Path not found for review move: {path}")
        return False

    try:
        review_root = Path(REVIEW_FOLDER)
        review_root.mkdir(parents=True, exist_ok=True)

        if path_obj.is_absolute():
            try:
                rel_path = path_obj.relative_to(Path.cwd())
            except Exception:
                rel_path = Path(path_obj.drive.rstrip(":")) / path_obj.relative_to(path_obj.anchor)
        else:
            rel_path = path_obj

        destination = review_root / rel_path
        destination.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(path_obj), str(destination))
        logging.warning(f"Moved file to review: {path_obj} -> {destination}")
        return True
    except Exception as exc:
        logging.error(f"Unable to move {path} to review folder: {exc}")
        return False


class _SuppressFd2:
    def __enter__(self):
        self._devnull = os.open(os.devnull, os.O_WRONLY)
        self._saved = os.dup(2)
        os.dup2(self._devnull, 2)
        return self

    def __exit__(self, *_):
        os.dup2(self._saved, 2)
        os.close(self._saved)
        os.close(self._devnull)


def _try_import(module, pip_name=None):
    try:
        import importlib
        return importlib.import_module(module)
    except ImportError:
        logging.warning(f"Optional dependency not installed: {pip_name or module}")
        return None

fitz_mod = _try_import("fitz", "pymupdf")
PIL_Image = _try_import("PIL.Image", "Pillow")
docx_mod = _try_import("docx", "python-docx")
openpyxl_mod = _try_import("openpyxl", "openpyxl")
pptx_mod = _try_import("pptx", "python-pptx")
langdetect_mod = _try_import("langdetect", "langdetect")
unstructured_mod = _try_import("unstructured.partition.auto", "unstructured")

if fitz_mod:
    try:
        import fitz
        with _SuppressFd2():
            fitz.TOOLS.mupdf_warnings()
            fitz.TOOLS.reset_mupdf_warnings()
    except Exception:
        pass

_TESSERACT_OK = False
_pytesseract = None
try:
    import pytesseract
    tess_exe = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    if os.path.exists(tess_exe):
        pytesseract.pytesseract.tesseract_cmd = tess_exe
    pytesseract.get_tesseract_version()
    _pytesseract = pytesseract
    _TESSERACT_OK = True
    print("  ✅ Tesseract found")
except Exception as _e:
    print(f"  ⚠  Tesseract disabled: {_e}")

_WIN32COM_OK = False
_WIN32COM_LOCK = threading.Lock()
try:
    import win32com.client
    _WIN32COM_OK = True
    print("  ✅ win32com found — .doc/.ppt enabled (serialized)")
except ImportError:
    print("  ⚠  win32com not found (pip install pywin32)")

PDF_EXT = {".pdf"}
DOCX_EXT = {".docx", ".doc", ".odt", ".rtf", ".wpd", ".wps", ".pages"}
XLSX_EXT = {".xlsx", ".xls", ".ods", ".xlsm", ".numbers"}
CSV_EXT = {".csv", ".tsv"}
PPTX_EXT = {".pptx", ".ppt", ".odp", ".ppsx", ".pps", ".key"}
ARCHIVE_EXT = {".zip", ".rar"}
JAR_EXT = {".jar", ".war"}
BINARY_TEXT_EXT = {".class", ".dex"}
TEXT_EXT = {
    ".txt", ".md", ".rst", ".nfo", ".diz", ".log",
    ".json", ".xml", ".yaml", ".yml", ".toml",
    ".ini", ".cfg", ".conf", ".env", ".properties",
    ".html", ".htm", ".css",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java",
    ".c", ".cpp", ".h", ".hpp", ".cs", ".go",
    ".rs", ".rb", ".sh", ".bash", ".csh", ".ps1", ".sql",
    ".r", ".tex", ".php", ".vb", ".lua", ".dart",
    ".kt", ".swift", ".pl", ".scala", ".m", ".asm",
    ".s", ".f", ".f90", ".jl", ".ex", ".exs",
    ".erl", ".hs", ".ml", ".clj", ".groovy", ".gradle",
    ".cmake", ".makefile", ".mk", ".bat", ".cmd",
    ".mf", ".jsp", ".jspx", ".xhtml", ".xsl", ".tld",
    ".pem", ".crt", ".cer", ".key", ".pub",
}
IMAGE_EXT = {
    ".jpg", ".jpeg", ".png", ".bmp",
    ".tiff", ".tif", ".gif", ".webp",
    ".ico", ".heic", ".heif",
}
SKIP_EXT = {
    ".exe", ".dll", ".sys", ".bin", ".dat",
    ".iso", ".vmdk", ".vhd", ".vhdx",
    ".mp3", ".mp4", ".avi", ".mov", ".mkv",
    ".flac", ".wav", ".aac", ".ogg",
    ".ttf", ".otf", ".woff", ".woff2",
    ".lnk", ".url", ".desktop",
    ".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".cab",
}
SKIP_CATEGORIES = {
    "System", "Executable", "Font", "VirtualMachine",
    "Shortcut", "Backup", "Temp", "Audio", "Video",
}

def _is_binary_file(path, check_bytes=8192, ext=None):
    try:
        if ext is None:
            ext = Path(path).suffix.lower()
        try:
            if ext in TEXT_EXT or ext in CSV_EXT:
                return False
        except Exception:
            pass
        with open(path, "rb") as f:
            chunk = f.read(check_bytes)
        if not chunk:
            return False
        text_chars = bytes(range(32, 127)) + b"\t\n\r\x0b\x0c"
        non_text = sum(1 for b in chunk if b not in text_chars)
        return (non_text / len(chunk)) > 0.30
    except Exception:
        return False


def extract_text_file(path):
    ext = Path(path).suffix.lower()
    if _is_binary_file(path, ext=ext):
        logging.warning(f"Binary file skipped in text extractor: {path}")
        return "", "FAILED_BINARY"

    for enc in ("utf-8", "latin-1", "cp1252", "cp1256"):
        try:
            text = Path(path).read_text(encoding=enc, errors="replace")
            return text.strip(), "plain_text_read"
        except Exception as e:
            logging.warning(f"Text read ({enc}) failed for {path}: {e}")
    return "", "FAILED"


def _extract_pdf_core(path):
    if not fitz_mod:
        return extract_via_unstructured(path)
    try:
        import fitz
        with _SuppressFd2():
            doc = fitz.open(path)
            fitz.TOOLS.mupdf_warnings()
            fitz.TOOLS.reset_mupdf_warnings()

        pages_text = []
        ocr_used = False
        for page in doc:
            with _SuppressFd2():
                try:
                    text = page.get_text("text").strip()
                    fitz.TOOLS.mupdf_warnings()
                except Exception:
                    text = ""

            if not text and _TESSERACT_OK and PIL_Image:
                try:
                    from PIL import Image
                    pix = page.get_pixmap(dpi=OCR_DPI)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    text = _pytesseract.image_to_string(img, lang=OCR_LANG).strip()
                    ocr_used = True
                except Exception as ocr_e:
                    logging.warning(f"Page OCR failed [{path}]: {ocr_e}")

            if text:
                pages_text.append(text)

        doc.close()
        full_text = "\n".join(pages_text).strip()
        if full_text:
            return full_text, ("pymupdf+ocr" if ocr_used else "pymupdf")
    except Exception as e:
        logging.warning(f"PyMuPDF failed for {path}: {e}")

    return extract_via_unstructured(path)


def extract_pdf(path):
    result = ["", "FAILED"]
    exc_info = [None]

    def _worker():
        try:
            result[0], result[1] = _extract_pdf_core(path)
        except Exception as e:
            exc_info[0] = e

    t = threading.Thread(target=_worker, daemon=True)
    t.start()
    t.join(timeout=PDF_TIMEOUT_SEC)

    if t.is_alive():
        logging.warning(f"PDF TIMEOUT (>{PDF_TIMEOUT_SEC}s): {path}")
        return "", "TIMEOUT"
    if exc_info[0]:
        logging.warning(f"PDF exception [{path}]: {exc_info[0]}")
        return "", "FAILED"

    return result[0], result[1]


def _doc_rtf_fallback(path):
    for enc in ("utf-8", "latin-1", "cp1252", "cp1256"):
        try:
            raw = Path(path).read_text(encoding=enc, errors="replace")
            if raw.startswith("{\\rtf"):
                text = re.sub(r"\\{\\*?\\\\[a-z]+[\\d\\-]* ?", " ", raw)
                text = re.sub(r"\\\\[a-z]+[\\d\\]* ?", " ", text)
                text = re.sub(r"[{}]", "", text)
                text = text.replace("\\par", "\n").replace("\\line", "\n")
                text = re.sub(r"\s+", " ", text).strip()
                if len(text) > 20:
                    return text, "rtf_fallback"
        except Exception:
            pass
    return "", "FAILED"


def extract_doc_win32(path):
    if not _WIN32COM_OK:
        return "", "FAILED"

    with _WIN32COM_LOCK:
        word = doc = None
        try:
            import win32com.client
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            word.DisplayAlerts = False
            doc = word.Documents.Open(str(Path(path).resolve()), ReadOnly=True)
            text = doc.Content.Text.strip()
            return text, "win32com_word"
        except Exception as e:
            logging.warning(f"win32com Word failed for {path}: {e}")
            return "", "FAILED"
        finally:
            try:
                if doc: doc.Close(False)
            except Exception:
                pass
            try:
                if word: word.Quit()
            except Exception:
                pass


def extract_ppt_win32(path):
    if not _WIN32COM_OK:
        return "", "FAILED"

    with _WIN32COM_LOCK:
        app = prs = None
        try:
            import win32com.client
            app = win32com.client.Dispatch("PowerPoint.Application")
            prs = app.Presentations.Open(
                str(Path(path).resolve()), ReadOnly=True, WithWindow=False
            )
            parts = []
            for slide in prs.Slides:
                for shape in slide.Shapes:
                    try:
                        if shape.HasTextFrame:
                            parts.append(shape.TextFrame.TextRange.Text.strip())
                    except Exception:
                        pass
            return "\n".join(p for p in parts if p), "win32com_ppt"
        except Exception as e:
            logging.warning(f"win32com PPT failed for {path}: {e}")
            return "", "FAILED"
        finally:
            try:
                if prs: prs.Close()
            except Exception:
                pass
            try:
                if app: app.Quit()
            except Exception:
                pass


def extract_docx_file(path):
    ext = path.lower()
    if ext.endswith(".docx") and docx_mod:
        try:
            from docx import Document
            doc = Document(path)
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
            text = "\n".join(parts).strip()
            if text:
                return text, "python-docx"
        except Exception as e:
            logging.warning(f"python-docx failed for {path}: {e}")
    if ext.endswith(".doc"):
        text, method = extract_doc_win32(path)
        if text:
            return text, method
        text, method = _doc_rtf_fallback(path)
        if text:
            return text, method
        return "", "FAILED"
    return extract_via_unstructured(path)


def extract_xlsx_file(path):
    if openpyxl_mod and path.lower().endswith((".xlsx", ".xlsm")):
        try:
            import openpyxl
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")
                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_text.strip():
                        parts.append(row_text)
            wb.close()
            text = "\n".join(parts).strip()
            if text:
                return text, "openpyxl"
        except Exception as e:
            logging.warning(f"openpyxl failed for {path}: {e}")
    return extract_via_unstructured(path)


def extract_pptx_file(path):
    ext = path.lower()
    if ext.endswith((".pptx", ".ppsx")) and pptx_mod:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                parts.append(t)
            text = "\n".join(parts).strip()
            if text:
                return text, "python-pptx"
        except Exception as e:
            logging.warning(f"python-pptx failed for {path}: {e}")
    if ext.endswith(".ppt"):
        text, method = extract_ppt_win32(path)
        if text:
            return text, method
    return extract_via_unstructured(path)


def extract_image(path):
    if not _TESSERACT_OK or not PIL_Image:
        return "", "SKIPPED_NO_OCR"
    try:
        from PIL import Image
        img = Image.open(path)
        text = _pytesseract.image_to_string(img, lang=OCR_LANG).strip()
        return text, "tesseract_ocr"
    except Exception as e:
        logging.warning(f"OCR failed for {path}: {e}")
    return "", "FAILED"


def extract_via_unstructured(path):
    if unstructured_mod:
        try:
            from unstructured.partition.auto import partition
            old_stderr = sys.stderr
            sys.stderr = io.StringIO()
            try:
                elements = partition(filename=path)
            finally:
                sys.stderr = old_stderr

            text = "\n".join(
                el.text for el in elements
                if hasattr(el, "text") and el.text
            ).strip()
            if text:
                return text, "unstructured"
        except Exception as e:
            if "soffice" not in str(e).lower():
                logging.warning(f"Unstructured failed for {path}: {e}")
    return "", "FAILED"


def _extract_bytes_to_temp(data, inner_name, category=""):
    ext = Path(inner_name).suffix.lower() or ".bin"
    text = ""
    method = "FAILED"
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        text, method = route_extraction(tmp_path, ext, category)
    except Exception as e:
        logging.warning(f"Temp extraction failed [{inner_name}]: {e}")
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
    return text, method, "ocr" in method.lower()


def extract_from_zip_entry(virtual_path, category=""):
    archive_path, inner_name = virtual_path.split("::", 1)
    try:
        with zipfile.ZipFile(archive_path, "r") as zf:
            data = zf.read(inner_name)
        return _extract_bytes_to_temp(data, inner_name, category)
    except Exception as e:
        logging.warning(f"ZIP failed [{inner_name}]: {e}")
        return "", "FAILED", False


def extract_from_rar_entry(virtual_path, category=""):
    if rarfile is None:
        logging.warning("RAR extraction requested but rarfile module not installed")
        return "", "SKIPPED_NO_RAR", False
    archive_path, inner_name = virtual_path.split("::", 1)
    try:
        with rarfile.RarFile(archive_path, "r") as rf:
            data = rf.read(inner_name)
        return _extract_bytes_to_temp(data, inner_name, category)
    except Exception as e:
        logging.warning(f"RAR failed [{inner_name}]: {e}")
        return "", "FAILED", False


def extract_jar_file(path, category=""):
    try:
        with zipfile.ZipFile(path, "r") as zf:
            allowed_exts = TEXT_EXT | CSV_EXT | PDF_EXT | DOCX_EXT | XLSX_EXT | PPTX_EXT | IMAGE_EXT | JAR_EXT | ARCHIVE_EXT | BINARY_TEXT_EXT
            parts = []
            for info in zf.infolist():
                if info.is_dir():
                    continue
                inner_ext = Path(info.filename).suffix.lower()
                if inner_ext not in allowed_exts:
                    continue
                try:
                    data = zf.read(info.filename)
                except Exception as e:
                    logging.warning(f"JAR entry read failed [{info.filename} in {path}]: {e}")
                    continue
                text, method, _ = _extract_bytes_to_temp(data, info.filename, category)
                if text:
                    parts.append(f"[{info.filename}]\n{text}")
                if len(parts) >= 20:
                    break
            full_text = "\n\n".join(parts).strip()
            if full_text:
                return full_text, "jar_archive"
    except Exception as e:
        logging.warning(f"JAR/WAR failed for {path}: {e}")
    return "", "FAILED"


def extract_archive_file(path, category=""):
    try:
        if path.lower().endswith(".zip"):
            with zipfile.ZipFile(path, "r") as zf:
                return _extract_archive_members(zf.read, zf.infolist(), category, path)
        if path.lower().endswith(".rar") and rarfile is not None:
            with rarfile.RarFile(path, "r") as rf:
                return _extract_archive_members(rf.read, rf.infolist(), category, path)
    except Exception as e:
        logging.warning(f"Archive extraction failed [{path}]: {e}")
    return "", "FAILED"


def _extract_archive_members(read_fn, members, category, archive_path):
    parts = []
    for info in members:
        name = getattr(info, 'filename', None) or getattr(info, 'name', None)
        if not name:
            continue
        if name.endswith('/') or name.endswith('\\'):
            continue
        inner_ext = Path(name).suffix.lower()
        if inner_ext not in (TEXT_EXT | CSV_EXT | PDF_EXT | DOCX_EXT | XLSX_EXT | PPTX_EXT | IMAGE_EXT | JAR_EXT | ARCHIVE_EXT | BINARY_TEXT_EXT):
            continue
        try:
            data = read_fn(name)
        except Exception as e:
            logging.warning(f"Archive member read failed [{name} in {archive_path}]: {e}")
            continue
        text, method, _ = _extract_bytes_to_temp(data, name, category)
        if text:
            parts.append(f"[{name}]\n{text}")
        if len(parts) >= 20:
            break
    full_text = "\n\n".join(parts).strip()
    return (full_text, "archive") if full_text else ("", "FAILED")


def extract_strings(path, extension, category=""):
    try:
        with open(path, "rb") as f:
            data = f.read()
    except Exception as e:
        logging.warning(f"String extraction failed for {path}: {e}")
        return "", "FAILED"
    min_len = 4
    max_chars = 100_000
    strings = []
    current = []
    for b in data:
        if 32 <= b < 127:
            current.append(chr(b))
        else:
            if len(current) >= min_len:
                strings.append(''.join(current))
            current = []
            if len(strings) >= 500:
                break
    if len(current) >= min_len:
        strings.append(''.join(current))
    text = "\n".join(strings)
    if text:
        return text[:max_chars], "strings_extracted"
    return "", "FAILED"


def route_extraction(path, extension, category=""):
    ext = extension.lower()
    if category in SKIP_CATEGORIES:
        return "", "SKIPPED"
    if ext in TEXT_EXT:
        return extract_text_file(path)
    if ext in CSV_EXT:
        return extract_text_file(path)
    if ext in PDF_EXT:
        return extract_pdf(path)
    if ext in DOCX_EXT:
        return extract_docx_file(path)
    if ext in XLSX_EXT:
        return extract_xlsx_file(path)
    if ext in PPTX_EXT:
        return extract_pptx_file(path)
    if ext in IMAGE_EXT:
        return extract_image(path)
    if ext in JAR_EXT:
        return extract_jar_file(path, category)
    if ext in ARCHIVE_EXT:
        return extract_archive_file(path, category)
    if ext in BINARY_TEXT_EXT:
        return extract_strings(path, ext, category)
    if ext in SKIP_EXT:
        return "", "SKIPPED"
    if category in ("Code", "Config", "Script", "Security", "Data"):
        return extract_text_file(path)
    return extract_via_unstructured(path)


def _detect_language(text):
    if not langdetect_mod:
        return ""
    if len(text) < MIN_LANG_CHARS:
        return ""
    try:
        from langdetect import detect
        return detect(text)
    except Exception:
        return ""


def classify_status(text, method, extension):
    if method in (
        "SKIPPED", "SKIPPED_EMPTY", "SKIPPED_NO_OCR", "SKIPPED_NO_RAR",
        "TIMEOUT", "NOT_REACHED",
    ):
        return "SKIPPED"
    if not text:
        return "FAILED"
    if extension in IMAGE_EXT and len(text) < 20:
        return "PARTIAL"
    if extension in (DOCX_EXT | PDF_EXT | XLSX_EXT | PPTX_EXT) and len(text) < 10:
        return "PARTIAL"
    return "SUCCESS"


def _safe_str(value) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    try:
        if pd.isna(value):
            return ""
    except Exception:
        pass
    return str(value)


def _process_single_row(args):
    loop_i, row = args

    path = str(row.get("path", ""))
    extension = str(row.get("extension", "")).lower()
    category = str(row.get("category", ""))
    is_empty = bool(row.get("is_empty", False))

    source_raw = row.get("source_archive", "")
    source = "" if pd.isna(source_raw) else str(source_raw)

    row_dict = dict(row)

    if is_empty and source == "":
        _move_to_review(path)
        return loop_i, None

    try:
        ocr_applied = False
        if source == "":
            text, method = route_extraction(path, extension, category)
            ocr_applied = "ocr" in method.lower()
        elif source.endswith(".zip"):
            text, method, ocr_applied = extract_from_zip_entry(path, category)
        elif source.endswith(".rar"):
            text, method, ocr_applied = extract_from_rar_entry(path, category)
        else:
            text, method = "", "SKIPPED"
    except Exception as e:
        logging.error(f"Unhandled error for {path}: {e}")
        text, method, ocr_applied = "", "ERROR", False

    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS]
        method += "_truncated"

    status = classify_status(text, method, extension)
    lang = _detect_language(text)

    if status in ("FAILED", "SKIPPED"):
        _move_to_review(path)
        return loop_i, None

    if lang and lang not in ALLOWED_LANGUAGES:
        _move_to_review(path)
        return loop_i, None

    return loop_i, {
        **row_dict,
        "file_id": row["file_id"],
        "extracted_text": text,
        "char_count": len(text),
        "word_count": len(text.split()) if text else 0,
        "extraction_method": method,
        "extraction_status": status,
        "ocr_applied": ocr_applied,
        "language_hint": lang,
        "phase1_status": str(row.get("status", "")),
    }


def _clean_text_for_db(text):
    if not isinstance(text, str):
        return text
    text = text.replace("\x00", "")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    return text


def get_db_connection():
    if not DB_URL:
        raise RuntimeError("DB_URL not set. Please set DB_URL in your .env or environment before running Phase 2.")
    return psycopg2.connect(DB_URL)


def create_tables(conn):
    with conn.cursor() as cur:
        cur.execute(
            """
            CREATE TABLE IF NOT EXISTS extracted_content (
                id SERIAL PRIMARY KEY,
                file_id INTEGER NOT NULL UNIQUE REFERENCES files(id) ON DELETE CASCADE,
                extracted_text TEXT,
                char_count INTEGER,
                word_count INTEGER,
                extraction_method TEXT,
                extraction_status TEXT,
                ocr_applied BOOLEAN,
                language_hint TEXT,
                phase1_status TEXT
            );
            """
        )
    conn.commit()


def fetch_unprocessed_files(conn):
    with conn.cursor(cursor_factory=RealDictCursor) as cur:
        cur.execute(
            """
            SELECT f.id AS file_id, f.*
            FROM files f
            LEFT JOIN extracted_content e ON f.id = e.file_id
            WHERE e.file_id IS NULL
            ORDER BY CASE WHEN f.status = 'KEEP' THEN 0 ELSE 1 END, f.id
            """
        )
        return cur.fetchall()


def insert_extracted_rows(conn, rows):
    if not rows:
        return
    sql = """
        INSERT INTO extracted_content (
            file_id, extracted_text, char_count, word_count,
            extraction_method, extraction_status, ocr_applied,
            language_hint, phase1_status
        ) VALUES %s
        ON CONFLICT (file_id) DO NOTHING
    """
    with conn.cursor() as cur:
        execute_values(cur, sql, rows)
    conn.commit()


def run_phase2(inventory_csv="file_inventory.csv", output_csv="extracted_content.csv"):
    print(f"\n{'═'*60}")
    print(f"  Phase 2 — Content Extraction  (v7)")
    print(f"{'═'*60}")
    print(f"\n  Capabilities:")
    print(f"    Tesseract    : {'✅' if _TESSERACT_OK else '❌ — images will be SKIPPED'}")
    print(f"    win32com     : {'✅ serialized (thread-safe)' if _WIN32COM_OK else '❌ (pip install pywin32)'}")
    print(f"    Unstructured : {'✅' if unstructured_mod else '❌'}")
    print(f"    Workers      : {NUM_WORKERS} parallel threads")
    print(f"    PDF timeout  : {PDF_TIMEOUT_SEC}s per file")
    print(f"    OCR DPI      : {OCR_DPI}")
    print(f"    Binary guard : ✅ (stops garbled reads on binary files)")
    print(f"    MuPDF noise  : ✅ suppressed (OS fd-2 redirect)")
    print(f"    Checkpoints  : ❌ disabled — DB resumes per file")
    print(f"    AI           : Phase 3\n")

    conn = get_db_connection()
    create_tables(conn)
    rows = fetch_unprocessed_files(conn)
    total = len(rows)
    print(f"  {total:,} entries loaded from files table")

    if total == 0:
        print("  No new files to extract.")
        conn.close()
        return 0

    success = 0
    partial = 0
    failed = 0
    skipped = 0
    reviewed = 0
    done_count = 0
    lock = threading.Lock()

    to_insert = []

    def flush_inserted():
        nonlocal to_insert
        if not to_insert:
            return
        insert_extracted_rows(conn, to_insert)
        to_insert = []

    print(f"  Priority: {sum(1 for row in rows if row['status'] == 'KEEP'):,} KEEP files first\n")
    print(f"  🚀 Starting parallel extraction with {NUM_WORKERS} workers...\n")

    try:
        with ThreadPoolExecutor(max_workers=NUM_WORKERS) as executor:
            futures = {
                executor.submit(_process_single_row, (i, row)): i
                for i, row in enumerate(rows)
            }
            for future in as_completed(futures):
                try:
                    loop_i, result = future.result()
                    with lock:
                        if result is None:
                            reviewed += 1
                        else:
                            st = result.get("extraction_status", "")
                            if st == "SUCCESS":
                                success += 1
                            elif st == "PARTIAL":
                                partial += 1
                            elif st == "FAILED":
                                failed += 1
                            else:
                                skipped += 1
                            to_insert.append(
                                (
                                    result["file_id"],
                                    _clean_text_for_db(result.get("extracted_text", "")),
                                    result.get("char_count", 0),
                                    result.get("word_count", 0),
                                    result.get("extraction_method", ""),
                                    result.get("extraction_status", ""),
                                    bool(result.get("ocr_applied", False)),
                                    result.get("language_hint", ""),
                                    result.get("phase1_status", ""),
                                )
                            )
                        done_count += 1
                        _done = done_count
                except Exception as e:
                    logging.error(f"Future result error: {e}")
                    with lock:
                        done_count += 1
                        failed += 1
                        _done = done_count
                if len(to_insert) >= 50:
                    flush_inserted()
                if _done % 50 == 0 or _done == total:
                    pct = _done / total * 100
                    bar = "█" * int(pct // 5) + "░" * (20 - int(pct // 5))
                    print(
                        f"  [{bar}] {pct:5.1f}%  ✅{success} ⚠{partial} ❌{failed} ⏭{skipped}",
                        end="\r"
                    )
        print()
    except KeyboardInterrupt:
        print(f"\n\n  ⚠  Interrupted — saving {done_count} rows collected so far...")
    except Exception as e:
        print(f"\n\n  ❌ Pipeline error: {e}")
        logging.error(f"Pipeline crash: {e}", exc_info=True)

    flush_inserted()
    conn.close()

    n = success + partial + failed + skipped
    print(f"\n{'═'*60}")
    print(f"  Extraction Summary")
    print(f"{'═'*60}")
    print(f"  Total        : {n:,}")
    if n > 0:
        print(f"  ✅ SUCCESS   : {success:,}  ({success/n*100:.1f}%)")
        print(f"  ⚠  PARTIAL   : {partial:,}  ({partial/n*100:.1f}%)")
        print(f"  ❌ FAILED    : {failed:,}  ({failed/n*100:.1f}%)")
        print(f"  ⏭ SKIPPED   : {skipped:,}  ({skipped/n*100:.1f}%)")
    else:
        print(f"  ✅ SUCCESS   : {success:,}  (0.0%)")
        print(f"  ⚠  PARTIAL   : {partial:,}  (0.0%)")
        print(f"  ❌ FAILED    : {failed:,}  (0.0%)")
        print(f"  ⏭ SKIPPED   : {skipped:,}  (0.0%)")
    print(f"  📁 Moved to review: {reviewed:,}")
    print(f"\n  Output → extracted_content table")
    print(f"  ➡  Ready for Phase 3 (Cleaning)")
    print(f"{'═'*60}")
    print("\nPhase 2 complete ✅")
    return n


if __name__ == "__main__":
    run_phase2()
